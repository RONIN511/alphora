"""
PPT 演示文稿查看器 - 处理 .pptx/.ppt 文件
"""
from typing import Optional, List, Tuple

from ..utils.common import get_file_info, truncate_text


class PresentationViewer:
    """PPT 演示文稿查看器"""
    
    SUPPORTED_EXTENSIONS = {'.pptx', '.ppt'}
    
    def __init__(self, file_path: str):
        self.file_path = file_path
        self.file_info = get_file_info(file_path)
        self.ext = self.file_info['extension']
        
    def view(
        self,
        purpose: str = "preview",
        keyword: Optional[str] = None,
        max_lines: int = 100,
        page_number: Optional[int] = None,
    ) -> str:
        """查看 PPT 内容"""
        purpose, warnings = self._infer_params(purpose, keyword)
        
        try:
            from pptx import Presentation
        except ImportError:
            return "❌ 需要安装 python-pptx: pip install python-pptx"
        
        try:
            prs = Presentation(self.file_path)
        except Exception as e:
            return f"❌ 无法打开演示文稿: {e}"
        
        total_slides = len(prs.slides)
        
        if purpose == "structure":
            return self._get_structure(prs, total_slides, warnings)
        elif purpose == "search":
            return self._search(prs, keyword, max_lines, warnings)
        elif page_number is not None:
            return self._view_slide(prs, page_number, total_slides, warnings)
        else:
            return self._preview(prs, total_slides, max_lines, warnings)
    
    def _infer_params(self, purpose: str, keyword: Optional[str]) -> Tuple[str, List[str]]:
        warnings = []
        if keyword and purpose != "search":
            warnings.append(f"⚠️ 检测到 keyword，已切换为 search 模式")
            purpose = "search"
        if purpose == "search" and not keyword:
            purpose = "preview"
        return purpose, warnings
    
    def _get_slide_title(self, slide) -> str:
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                return truncate_text(shape.text.strip(), 50)
        return "(无标题)"
    
    def _format_header(self, total_slides: int, warnings: List[str]) -> str:
        lines = [
            f"📊 文件: {self.file_info['name']}",
            f"📦 大小: {self.file_info['size_human']}",
            f"📋 幻灯片: {total_slides} 页",
        ]
        if warnings:
            lines.extend([""] + warnings)
        return '\n'.join(lines)
    
    def _get_structure(self, prs, total_slides, warnings) -> str:
        lines = [self._format_header(total_slides, warnings), "", "【幻灯片列表】"]
        for i, slide in enumerate(prs.slides, 1):
            title = self._get_slide_title(slide)
            lines.append(f"  第{i}页: {title}")
        return '\n'.join(lines)
    
    def _preview(self, prs, total_slides, max_lines, warnings) -> str:
        lines = [self._format_header(total_slides, warnings), "", "【内容预览】"]
        count = 0
        for i, slide in enumerate(prs.slides, 1):
            if count >= max_lines:
                break
            title = self._get_slide_title(slide)
            lines.append(f"\n━━━ 第{i}页: {title} ━━━")
            for shape in slide.shapes:
                if count >= max_lines:
                    break
                if hasattr(shape, 'text') and shape.text.strip():
                    lines.append(shape.text.strip()[:200])
                    count += 1
        return '\n'.join(lines)
    
    def _view_slide(self, prs, page_number, total_slides, warnings) -> str:
        if page_number < 1 or page_number > total_slides:
            return f"❌ 页码超出范围 (1-{total_slides})"
        
        slide = prs.slides[page_number - 1]
        lines = [
            f"📊 文件: {self.file_info['name']}",
            f"📋 第 {page_number}/{total_slides} 页", "",
            "【页面内容】"
        ]
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                lines.append(shape.text.strip())
        return '\n'.join(lines)
    
    def _search(self, prs, keyword, max_lines, warnings) -> str:
        results = []
        keyword_lower = keyword.lower()
        
        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, 'text') and keyword_lower in shape.text.lower():
                    results.append({'page': i, 'content': shape.text.strip()[:100]})
        
        lines = [
            self._format_header(len(prs.slides), warnings), "",
            f"🔍 搜索: '{keyword}'",
            f"📋 找到 {len(results)} 处匹配", ""
        ]
        
        if results:
            for r in results[:max_lines]:
                lines.append(f"[第{r['page']}页] {r['content']}")
        else:
            lines.append(f"未找到包含 '{keyword}' 的内容")
        
        return '\n'.join(lines)
